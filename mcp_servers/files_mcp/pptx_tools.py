"""Ferramentas PowerPoint: criar e editar apresentações personalizadas."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_DIR = Path(__file__).parent.parent.parent / "outputs" / "presentations"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Paleta padrão Clilink
BLUE = RGBColor(0x00, 0x77, 0xB5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xFA, 0xA3, 0x2C)


def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text: str, left, top, width, height, font_size: int, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def create_presentation(title: str, slides_data: list[dict], filename: str = None) -> str:
    """
    Cria apresentação PowerPoint.
    slides_data = [
      {"type": "title", "title": "...", "subtitle": "..."},
      {"type": "content", "title": "...", "bullets": ["ponto 1", "ponto 2"]},
      {"type": "section", "title": "..."},
      {"type": "closing", "title": "...", "subtitle": "..."},
    ]
    """
    fname = filename or f"{title.lower().replace(' ', '_')}.pptx"
    path = OUTPUT_DIR / fname
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_data.get("type", "content")

        if stype == "title":
            _set_bg(slide, DARK)
            # Faixa azul lateral
            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
            shape.line.fill.background()
            _add_text(slide, slide_data.get("title", ""), Inches(0.8), Inches(2.5), Inches(11), Inches(1.5), 40, bold=True, align=PP_ALIGN.LEFT)
            _add_text(slide, slide_data.get("subtitle", ""), Inches(0.8), Inches(4.2), Inches(11), Inches(1), 20, color=ACCENT, align=PP_ALIGN.LEFT)

        elif stype == "section":
            _set_bg(slide, BLUE)
            _add_text(slide, slide_data.get("title", ""), Inches(1), Inches(2.8), Inches(11), Inches(2), 36, bold=True, align=PP_ALIGN.CENTER)

        elif stype == "content":
            _set_bg(slide, WHITE)
            # Header azul
            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
            shape.line.fill.background()
            _add_text(slide, slide_data.get("title", ""), Inches(0.4), Inches(0.1), Inches(12), Inches(1), 24, bold=True, align=PP_ALIGN.LEFT)
            bullets = slide_data.get("bullets", [])
            y = Inches(1.5)
            for bullet in bullets[:8]:
                _add_text(slide, f"▸  {bullet}", Inches(0.6), y, Inches(12), Inches(0.6), 16, color=DARK)
                y += Inches(0.65)

        elif stype == "closing":
            _set_bg(slide, DARK)
            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT
            shape.line.fill.background()
            _add_text(slide, slide_data.get("title", "Obrigado!"), Inches(0.8), Inches(2.8), Inches(11), Inches(1.5), 40, bold=True, align=PP_ALIGN.LEFT)
            _add_text(slide, slide_data.get("subtitle", ""), Inches(0.8), Inches(4.5), Inches(11), Inches(1), 18, color=ACCENT, align=PP_ALIGN.LEFT)

    prs.save(str(path))
    return str(path)


def add_slide(pptx_path: str, slide_data: dict) -> str:
    """Adiciona um slide a uma apresentação existente."""
    prs = Presentation(pptx_path)
    # Reutiliza create_presentation para o slide
    tmp_path = create_presentation("tmp", [slide_data], "tmp_slide.pptx")
    tmp_prs = Presentation(tmp_path)
    template = tmp_prs.slides[0]
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)
    for shape in template.shapes:
        el = shape.element
        new_slide.shapes._spTree.insert(2, el)
    prs.save(pptx_path)
    Path(tmp_path).unlink(missing_ok=True)
    return f"Slide adicionado em: {pptx_path}"
