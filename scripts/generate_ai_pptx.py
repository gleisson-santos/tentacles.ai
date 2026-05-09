
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configuração de caminhos
BASE_DIR = Path(__file__).parent.parent
OUTPUT_DIR = BASE_DIR / "outputs" / "presentations"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
FILENAME = OUTPUT_DIR / "impacto_ia_corporativo.pptx"

# Cores (Paleta Clilink/Dark Modern)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x77, 0xB5)
ACCENT = RGBColor(0xFA, 0xA3, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Título ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, DARK)
    # Faixa decorativa
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    
    add_text(slide, "O IMPACTO DAS IAs NO\nMUNDO CORPORATIVO", Inches(1), Inches(2.2), Inches(11), Inches(2), 44, bold=True)
    add_text(slide, "Transformação Digital e a Nova Era da Eficiência", Inches(1), Inches(4.5), Inches(11), Inches(1), 24, color=ACCENT)

    # --- SLIDE 2: Índice ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    # Header
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()
    add_text(slide, "Conteúdo da Apresentação", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), 32, bold=True)
    
    topics = [
        "1. Automação de Processos Repetitivos",
        "2. Análise de Dados e Tomada de Decisão",
        "3. Personalização da Experiência do Cliente",
        "4. Otimização da Cadeia de Suprimentos",
        "5. O Futuro do Trabalho e Novas Habilidades"
    ]
    y = Inches(1.8)
    for topic in topics:
        add_text(slide, f"• {topic}", Inches(1), y, Inches(11), Inches(0.7), 22, color=DARK)
        y += Inches(0.8)

    # --- SLIDE 3: Automação ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()
    add_text(slide, "1. Automação Inteligente", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), 32, bold=True)
    
    points = [
        "Eliminação de tarefas manuais de baixo valor agregado.",
        "RPA (Robotic Process Automation) integrado com IA Generativa.",
        "Redução drástica de erros operacionais e aumento de velocidade.",
        "Exemplo: Processamento automático de faturas e contratos."
    ]
    y = Inches(1.8)
    for p in points:
        add_text(slide, f"▸ {p}", Inches(1), y, Inches(11), Inches(0.7), 20, color=DARK)
        y += Inches(0.9)

    # --- SLIDE 4: Decisões Baseadas em Dados ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()
    add_text(slide, "2. Análise Preditiva e Decisões", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), 32, bold=True)
    
    points = [
        "Processamento de Big Data em tempo real para insights estratégicos.",
        "Modelos preditivos que antecipam tendências de mercado.",
        "Democratização da análise de dados via linguagem natural (LLMs).",
        "Decisões baseadas em evidências, não apenas em intuição."
    ]
    y = Inches(1.8)
    for p in points:
        add_text(slide, f"▸ {p}", Inches(1), y, Inches(11), Inches(0.7), 20, color=DARK)
        y += Inches(0.9)

    # --- SLIDE 5: Futuro do Trabalho ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()
    add_text(slide, "5. O Futuro do Trabalho", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), 32, bold=True)
    
    points = [
        "IA como copiloto, não como substituto integral.",
        "Necessidade de 'Reskilling' e 'Upskilling' da força de trabalho.",
        "Foco em habilidades humanas: Criatividade, Empatia e Estratégia.",
        "Colaboração Humano-IA: Aumentando o potencial individual."
    ]
    y = Inches(1.8)
    for p in points:
        add_text(slide, f"▸ {p}", Inches(1), y, Inches(11), Inches(0.7), 20, color=DARK)
        y += Inches(0.9)

    # --- SLIDE 6: Fechamento ---
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, DARK)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    
    add_text(slide, "ESTEJA PREPARADO PARA A\nREVOLUÇÃO DA IA", Inches(1), Inches(2.5), Inches(11), Inches(2), 36, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Obrigado! • Perguntas e Respostas", Inches(1), Inches(4.8), Inches(11), Inches(1), 20, color=ACCENT, align=PP_ALIGN.CENTER)

    prs.save(str(FILENAME))
    print(f"Apresentação PPTX gerada com sucesso em: {FILENAME}")

if __name__ == "__main__":
    generate_pptx()
